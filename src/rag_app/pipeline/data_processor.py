"""
Professional Document Processor for RAG System
Handles document upload, text extraction, chunking, and embedding generation
"""

import os
import logging
from typing import List, Dict, Any, Optional, Tuple
import pickle
import numpy as np
from django.conf import settings
from django.utils import timezone
from ..models import Document as DocumentModel, DocumentChunk

# Import packages with proper error handling
try:
    import fitz  # PyMuPDF for better PDF handling
    PYMUPDF_AVAILABLE = True
except ImportError:
    PYMUPDF_AVAILABLE = False
    
try:
    from PyPDF2 import PdfReader
    PYPDF2_AVAILABLE = True
except ImportError:
    PYPDF2_AVAILABLE = False

try:
    from docx import Document as DocxDocument
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    from langchain_text_splitters import RecursiveCharacterTextSplitter
    from langchain_core.documents import Document
    LANGCHAIN_AVAILABLE = True
except ImportError:
    try:
        # Fallback to older langchain imports
        from langchain.text_splitter import RecursiveCharacterTextSplitter
        from langchain.schema import Document
        LANGCHAIN_AVAILABLE = True
    except ImportError:
        LANGCHAIN_AVAILABLE = False

try:
    from sentence_transformers import SentenceTransformer
    SENTENCE_TRANSFORMERS_AVAILABLE = True
except ImportError:
    SENTENCE_TRANSFORMERS_AVAILABLE = False

logger = logging.getLogger(__name__)


class DocumentProcessingError(Exception):
    """Custom exception for document processing errors"""
    pass


class DocumentProcessor:
    """
    Professional document processor for RAG system
    
    Features:
    - Supports multiple file formats (PDF, DOCX, TXT, PPTX)
    - Intelligent text chunking with LangChain
    - Embedding generation with SentenceTransformers
    - Database integration with Django models
    - Comprehensive error handling and logging
    """
    
    def __init__(self, 
                 chunk_size: int = 1000,
                 chunk_overlap: int = 200,
                 embedding_model: str = 'all-MiniLM-L6-v2'):
        """
        Initialize the document processor
        
        Args:
            chunk_size: Maximum size of text chunks
            chunk_overlap: Overlap between consecutive chunks
            embedding_model: Name of the sentence transformer model
        """
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.embedding_model_name = embedding_model
        
        # Initialize text splitter
        if LANGCHAIN_AVAILABLE:
            self.text_splitter = RecursiveCharacterTextSplitter(
                chunk_size=chunk_size,
                chunk_overlap=chunk_overlap,
                length_function=len,
                separators=["\n\n", "\n", ".", "!", "?", " ", ""]
            )
        else:
            logger.warning("LangChain not available, using basic text splitting")
            self.text_splitter = None
        
        # Initialize embedding model
        if SENTENCE_TRANSFORMERS_AVAILABLE:
            try:
                self.embedding_model = SentenceTransformer(embedding_model)
                logger.info(f"Loaded embedding model: {embedding_model}")
            except Exception as e:
                logger.error(f"Failed to load embedding model: {e}")
                self.embedding_model = None
        else:
            logger.warning("SentenceTransformers not available, embeddings disabled")
            self.embedding_model = None
    
    def process_document(self, document: DocumentModel) -> Dict[str, Any]:
        """
        Process a document and create chunks with embeddings
        
        Args:
            document: Django Document model instance
            
        Returns:
            Dict with processing results and statistics
        """
        start_time = timezone.now()
        
        try:
            logger.info(f"Starting processing for document {document.id}: {document.title}")
            
            # Validate document
            if not os.path.exists(document.file.path):
                raise DocumentProcessingError(f"File not found: {document.file.path}")
            
            # Extract text based on file type
            text_content, metadata = self._extract_text_with_metadata(document)
            
            if not text_content or not text_content.strip():
                raise DocumentProcessingError("No text content extracted from document")
            
            logger.info(f"Extracted {len(text_content)} characters from {document.title}")
            
            # Create chunks
            chunks = self._create_chunks(text_content, document, metadata)
            
            if not chunks:
                raise DocumentProcessingError("No chunks created from document")
            
            logger.info(f"Created {len(chunks)} chunks from {document.title}")
            
            # Generate embeddings and save chunks
            saved_chunks = self._create_embeddings_and_save(chunks, document)
            
            # Update document metadata
            document.processed = True
            document.processed_at = timezone.now()
            document.page_count = metadata.get('page_count', 0)
            document.save()
            
            processing_time = (timezone.now() - start_time).total_seconds()
            
            result = {
                'success': True,
                'document_id': str(document.id),
                'chunks_created': len(saved_chunks),
                'total_characters': len(text_content),
                'processing_time': processing_time,
                'page_count': metadata.get('page_count', 0),
                'metadata': metadata
            }
            
            logger.info(f"Successfully processed document {document.id} in {processing_time:.2f}s")
            return result
            
        except Exception as e:
            error_msg = f"Error processing document {document.id}: {str(e)}"
            logger.error(error_msg)
            
            return {
                'success': False,
                'document_id': str(document.id),
                'error': str(e),
                'processing_time': (timezone.now() - start_time).total_seconds()
            }
    
    def process_temp_document(self, temp_doc) -> Dict[str, Any]:
        """
        Process a temporary document for anonymous chat
        
        Args:
            temp_doc: TempDocument instance
            
        Returns:
            Dict with processing results
        """
        start_time = timezone.now()
        
        try:
            logger.info(f"Processing temporary document: {temp_doc.title}")
            
            # Validate temp document file exists
            if not os.path.exists(temp_doc.file.path):
                raise DocumentProcessingError(f"Temp file not found: {temp_doc.file.path}")
            
            # Extract text content from the temporary document
            text_content = self._extract_temp_document_text(temp_doc)
            
            if not text_content or not text_content.strip():
                raise DocumentProcessingError("No text content extracted from temporary document")
            
            logger.info(f"Extracted {len(text_content)} characters from temp document {temp_doc.title}")
            
            # Store the extracted text content in a cache for later retrieval
            # For simplicity, we'll store it as a file attribute or in cache
            cache_key = f"temp_doc_content_{temp_doc.id}"
            from django.core.cache import cache
            cache.set(cache_key, text_content, timeout=86400)  # 24 hours
            
            temp_doc.processed = True
            temp_doc.save()
            
            processing_time = (timezone.now() - start_time).total_seconds()
            
            logger.info(f"Successfully processed temp document {temp_doc.id} in {processing_time:.2f}s")
            
            return {
                'success': True,
                'temp_document_id': str(temp_doc.id),
                'processing_time': processing_time,
                'text_length': len(text_content),
                'message': 'Temporary document processed successfully'
            }
            
        except Exception as e:
            logger.error(f"Error processing temp document {temp_doc.id}: {str(e)}")
            return {
                'success': False,
                'temp_document_id': str(temp_doc.id),
                'error': str(e),
                'processing_time': (timezone.now() - start_time).total_seconds()
            }
    
    def _extract_text_with_metadata(self, document: DocumentModel) -> Tuple[str, Dict[str, Any]]:
        """
        Extract text content and metadata from uploaded file
        
        Returns:
            Tuple of (text_content, metadata_dict)
        """
        file_path = document.file.path
        file_extension = document.document_type.lower()
        metadata = {'file_type': file_extension, 'page_count': 0}
        
        try:
            if file_extension == 'pdf':
                text, page_count = self._extract_pdf_text(file_path)
                metadata['page_count'] = page_count
                return text, metadata
            elif file_extension == 'docx':
                return self._extract_docx_text(file_path), metadata
            elif file_extension == 'txt':
                return self._extract_txt_text(file_path), metadata
            elif file_extension == 'pptx':
                text, slide_count = self._extract_pptx_text(file_path)
                metadata['slide_count'] = slide_count
                return text, metadata
            else:
                raise DocumentProcessingError(f"Unsupported file type: {file_extension}")
                
        except Exception as e:
            raise DocumentProcessingError(f"Error extracting text from {file_path}: {str(e)}")
    
    def _extract_pdf_text(self, file_path: str) -> Tuple[str, int]:
        """Extract text from PDF file with page counting"""
        text = ""
        page_count = 0
        
        # Try PyMuPDF first (better quality)
        if PYMUPDF_AVAILABLE:
            try:
                doc = fitz.open(file_path)
                page_count = len(doc)
                for page_num in range(page_count):
                    page = doc.load_page(page_num)
                    page_text = page.get_text()
                    if page_text.strip():  # Only add non-empty pages
                        text += f"\n--- Page {page_num + 1} ---\n"
                        text += page_text
                doc.close()
                logger.info(f"Extracted text from {page_count} pages using PyMuPDF")
                return text, page_count
            except Exception as e:
                logger.warning(f"PyMuPDF extraction failed: {e}, trying PyPDF2")
        
        # Fallback to PyPDF2
        if PYPDF2_AVAILABLE:
            try:
                reader = PdfReader(file_path)
                page_count = len(reader.pages)
                for i, page in enumerate(reader.pages):
                    page_text = page.extract_text()
                    if page_text.strip():
                        text += f"\n--- Page {i + 1} ---\n"
                        text += page_text
                logger.info(f"Extracted text from {page_count} pages using PyPDF2")
                return text, page_count
            except Exception as e:
                logger.error(f"PyPDF2 extraction also failed: {e}")
        
        raise DocumentProcessingError("No PDF processing library available")
    
    def _extract_docx_text(self, file_path: str) -> str:
        """Extract text from DOCX file"""
        if not DOCX_AVAILABLE:
            raise DocumentProcessingError("python-docx not available")
        
        try:
            doc = DocxDocument(file_path)
            text_parts = []
            
            # Extract paragraphs
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_parts.append(paragraph.text)
            
            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        text_parts.append(" | ".join(row_text))
            
            text = "\n".join(text_parts)
            logger.info(f"Extracted {len(text_parts)} paragraphs/elements from DOCX")
            return text
            
        except Exception as e:
            raise DocumentProcessingError(f"Error extracting DOCX text: {str(e)}")
    
    def _extract_txt_text(self, file_path: str) -> str:
        """Extract text from TXT file with encoding detection"""
        encodings = ['utf-8', 'utf-16', 'latin-1', 'cp1252']
        
        for encoding in encodings:
            try:
                with open(file_path, 'r', encoding=encoding) as file:
                    text = file.read()
                    logger.info(f"Successfully read TXT file with {encoding} encoding")
                    return text
            except UnicodeDecodeError:
                continue
            except Exception as e:
                logger.error(f"Error reading TXT file with {encoding}: {e}")
                continue
        
        raise DocumentProcessingError("Unable to read TXT file with any encoding")
    
    def _extract_pptx_text(self, file_path: str) -> Tuple[str, int]:
        """Extract text from PPTX file with slide counting"""
        if not PPTX_AVAILABLE:
            raise DocumentProcessingError("python-pptx not available")
        
        try:
            prs = Presentation(file_path)
            text_parts = []
            slide_count = len(prs.slides)
            
            for i, slide in enumerate(prs.slides):
                slide_text = f"\n--- Slide {i + 1} ---\n"
                slide_content = []
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_content.append(shape.text.strip())
                
                if slide_content:
                    slide_text += "\n".join(slide_content)
                    text_parts.append(slide_text)
            
            text = "\n".join(text_parts)
            logger.info(f"Extracted text from {slide_count} slides")
            return text, slide_count
            
        except Exception as e:
            raise DocumentProcessingError(f"Error extracting PPTX text: {str(e)}")
    
    def _create_chunks(self, text: str, document: DocumentModel, metadata: Dict[str, Any]) -> List[Document]:
        """
        Split text into chunks using LangChain or basic splitting
        """
        base_metadata = {
            'source': document.title,
            'document_id': str(document.id),
            'subject_id': str(document.subject.id) if document.subject else None,
            'subject_name': document.subject.name if document.subject else None,
            'uploaded_by': document.uploaded_by.username,
            'file_type': metadata.get('file_type'),
            'page_count': metadata.get('page_count', 0)
        }
        
        if LANGCHAIN_AVAILABLE and self.text_splitter:
            # Use LangChain for intelligent splitting
            doc = Document(page_content=text, metadata=base_metadata)
            chunks = self.text_splitter.split_documents([doc])
            logger.info(f"Created {len(chunks)} chunks using LangChain")
        else:
            # Basic splitting fallback
            chunks = self._basic_text_split(text, base_metadata)
            logger.info(f"Created {len(chunks)} chunks using basic splitting")
        
        return chunks
    
    def _basic_text_split(self, text: str, metadata: Dict[str, Any]) -> List[Document]:
        """
        Basic text splitting when LangChain is not available
        """
        chunks = []
        chunk_size = self.chunk_size
        overlap = self.chunk_overlap
        
        # Split by paragraphs first, then by sentences
        paragraphs = text.split('\n\n')
        current_chunk = ""
        
        for paragraph in paragraphs:
            if len(current_chunk) + len(paragraph) <= chunk_size:
                current_chunk += paragraph + "\n\n"
            else:
                if current_chunk.strip():
                    chunks.append(Document(
                        page_content=current_chunk.strip(),
                        metadata=metadata.copy()
                    ))
                current_chunk = paragraph + "\n\n"
        
        # Add the last chunk
        if current_chunk.strip():
            chunks.append(Document(
                page_content=current_chunk.strip(),
                metadata=metadata.copy()
            ))
        
        return chunks
    
    def _create_embeddings_and_save(self, chunks: List[Document], document: DocumentModel) -> List[DocumentChunk]:
        """
        Create embeddings for chunks and save to database
        """
        # Delete existing chunks for this document
        DocumentChunk.objects.filter(document=document).delete()
        
        saved_chunks = []
        
        for i, chunk in enumerate(chunks):
            try:
                # Create embedding if model is available
                embedding_bytes = None
                if self.embedding_model:
                    embedding = self.embedding_model.encode(chunk.page_content)
                    embedding_bytes = pickle.dumps(embedding.astype(np.float32))
                
                # Extract page number from content if available
                page_number = self._extract_page_number(chunk.page_content)
                
                # Create and save chunk
                doc_chunk = DocumentChunk.objects.create(
                    document=document,
                    content=chunk.page_content,
                    chunk_index=i,
                    page_number=page_number,
                    embedding_vector=embedding_bytes
                )
                
                saved_chunks.append(doc_chunk)
                
            except Exception as e:
                logger.error(f"Error creating chunk {i} for document {document.id}: {e}")
                continue
        
        logger.info(f"Saved {len(saved_chunks)} chunks to database")
        return saved_chunks
    
    def _extract_page_number(self, content: str) -> Optional[int]:
        """Extract page number from chunk content if present"""
        import re
        
        # Look for page markers like "--- Page 1 ---"
        page_match = re.search(r'--- Page (\d+) ---', content)
        if page_match:
            return int(page_match.group(1))
        
        return None
    
    def get_document_chunks(self, document: DocumentModel) -> List[DocumentChunk]:
        """Get all chunks for a document"""
        return DocumentChunk.objects.filter(document=document).order_by('chunk_index')
    
    def get_subject_chunks(self, subject_id: int) -> List[DocumentChunk]:
        """Get all chunks for documents in a subject"""
        return DocumentChunk.objects.filter(
            document__subject_id=subject_id,
            document__processed=True
        ).order_by('document__title', 'chunk_index')
    
    def get_processing_stats(self) -> Dict[str, Any]:
        """Get processing statistics"""
        from django.db.models import Count, Sum
        
        stats = {
            'total_documents': DocumentModel.objects.filter(processed=True).count(),
            'total_chunks': DocumentChunk.objects.count(),
            'documents_by_type': DocumentModel.objects.filter(processed=True).values('document_type').annotate(count=Count('id')),
            'chunks_by_document': DocumentChunk.objects.values('document__title').annotate(count=Count('id')).order_by('-count')[:10]
        }
        
        return stats
    
    def _extract_temp_document_text(self, temp_doc) -> str:
        """
        Extract text content from a temporary document
        
        Args:
            temp_doc: TempDocument instance
            
        Returns:
            str: Extracted text content
        """
        file_path = temp_doc.file.path
        
        # Determine file type from file extension
        file_extension = os.path.splitext(file_path)[1].lower()
        
        try:
            if file_extension == '.pdf':
                text, _ = self._extract_pdf_text(file_path)
                return text
            elif file_extension == '.docx':
                return self._extract_docx_text(file_path)
            elif file_extension == '.txt':
                return self._extract_txt_text(file_path)
            elif file_extension == '.pptx':
                return self._extract_pptx_text(file_path)
            else:
                logger.warning(f"Unsupported file type for temp document: {file_extension}")
                return f"Content of {temp_doc.title} (unsupported file type: {file_extension})"
                
        except Exception as e:
            logger.error(f"Error extracting text from temp document {temp_doc.id}: {e}")
            return f"Error reading content from {temp_doc.title}: {str(e)}"
