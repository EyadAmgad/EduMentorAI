from django.shortcuts import render, get_object_or_404, redirect
from django.contrib.auth.decorators import login_required
from django.contrib.auth.mixins import LoginRequiredMixin
from django.views.generic import (
    TemplateView, ListView, DetailView, CreateView, 
    UpdateView, DeleteView, FormView
)
from django.views import View
from django.contrib import messages
from django.http import JsonResponse, HttpResponse, Http404, HttpResponseForbidden, HttpResponseServerError
from django.urls import reverse_lazy, reverse
from django.utils import timezone
from django.db.models import Q, Count, Avg
from django.core.paginator import Paginator
from django.views.decorators.csrf import csrf_exempt
from django.utils.decorators import method_decorator
from django.conf import settings
from io import BytesIO
import json
import logging
import os

from .models import (
    Subject, Document, DocumentChunk, ChatSession, ChatMessage,
    Quiz, Question, QuizAttempt, QuizResponse, UserProfile, StudySession,
    TempDocument
)
from .forms import (
    SubjectForm, DocumentUploadForm, QuizCreateForm, UserProfileForm,
    ChatMessageForm
)
from .pipeline.data_processor import DocumentProcessor

logger = logging.getLogger(__name__)


class HomeView(TemplateView):
    """Landing page view"""
    template_name = 'rag_app/home.html'
    
    def get_context_data(self, **kwargs):
        context = super().get_context_data(**kwargs)
        if self.request.user.is_authenticated:
            context['recent_subjects'] = Subject.objects.filter(
                created_by=self.request.user
            ).order_by('-created_at')[:5]
        return context


class DashboardView(LoginRequiredMixin, TemplateView):
    """User dashboard view"""
    template_name = 'rag_app/dashboard.html'
    
    def get_context_data(self, **kwargs):
        context = super().get_context_data(**kwargs)
        user = self.request.user
        
        # Get user statistics
        context.update({
            'total_subjects': Subject.objects.filter(created_by=user).count(),
            'total_documents': Document.objects.filter(uploaded_by=user).count(),
            'total_quizzes': Quiz.objects.filter(created_by=user).count(),
            'recent_documents': Document.objects.filter(uploaded_by=user).order_by('-uploaded_at')[:5],
            'recent_quizzes': Quiz.objects.filter(created_by=user).order_by('-created_at')[:5],
            'recent_chat_sessions': ChatSession.objects.filter(user=user).order_by('-last_activity')[:5],
            'quiz_attempts': QuizAttempt.objects.filter(user=user, is_completed=True).order_by('-completed_at')[:5],
        })
        
        # Calculate average quiz score
        avg_score = QuizAttempt.objects.filter(
            user=user, is_completed=True
        ).aggregate(avg_score=Avg('score'))['avg_score']
        context['average_quiz_score'] = round(avg_score, 1) if avg_score else 0
        
        return context


# Subject Views
class SubjectListView(LoginRequiredMixin, ListView):
    """List all subjects for the user"""
    model = Subject
    template_name = 'rag_app/subject_list.html'
    context_object_name = 'subjects'
    paginate_by = 10
    
    def get_queryset(self):
        return Subject.objects.filter(created_by=self.request.user)


class SubjectDetailView(LoginRequiredMixin, DetailView):
    """Subject detail view"""
    model = Subject
    template_name = 'rag_app/subject_detail.html'
    context_object_name = 'subject'
    
    def get_queryset(self):
        return Subject.objects.filter(created_by=self.request.user)
    
    def get_context_data(self, **kwargs):
        context = super().get_context_data(**kwargs)
        subject = self.get_object()
        context.update({
            'documents': subject.documents.all(),
            'quizzes': subject.quizzes.all(),
            'document_count': subject.documents.count(),
            'quiz_count': subject.quizzes.count(),
        })
        return context


class SubjectCreateView(LoginRequiredMixin, CreateView):
    """Create new subject"""
    model = Subject
    form_class = SubjectForm
    template_name = 'rag_app/subject_form.html'
    success_url = reverse_lazy('rag_app:subject_list')
    
    def form_valid(self, form):
        form.instance.created_by = self.request.user
        messages.success(self.request, 'Subject created successfully!')
        return super().form_valid(form)


class SubjectUpdateView(LoginRequiredMixin, UpdateView):
    """Update existing subject"""
    model = Subject
    form_class = SubjectForm
    template_name = 'rag_app/subject_form.html'
    
    def get_queryset(self):
        return Subject.objects.filter(created_by=self.request.user)
    
    def get_success_url(self):
        return reverse('rag_app:subject_detail', kwargs={'pk': self.object.pk})
    
    def form_valid(self, form):
        messages.success(self.request, 'Subject updated successfully!')
        return super().form_valid(form)


class SubjectDeleteView(LoginRequiredMixin, DeleteView):
    """Delete subject"""
    model = Subject
    template_name = 'rag_app/subject_confirm_delete.html'
    success_url = reverse_lazy('rag_app:subject_list')
    
    def get_queryset(self):
        return Subject.objects.filter(created_by=self.request.user)
    
    def delete(self, request, *args, **kwargs):
        messages.success(request, 'Subject deleted successfully!')
        return super().delete(request, *args, **kwargs)


# Document Views
class DocumentListView(LoginRequiredMixin, ListView):
    """List all documents for the user"""
    model = Document
    template_name = 'rag_app/document_list.html'
    context_object_name = 'documents'
    paginate_by = 12
    
    def get_queryset(self):
        queryset = Document.objects.filter(uploaded_by=self.request.user)
        subject_id = self.request.GET.get('subject')
        if subject_id:
            queryset = queryset.filter(subject_id=subject_id)
        return queryset.order_by('-uploaded_at')
    
    def get_context_data(self, **kwargs):
        context = super().get_context_data(**kwargs)
        context['subjects'] = Subject.objects.filter(created_by=self.request.user)
        context['selected_subject'] = self.request.GET.get('subject')
        return context


class DocumentDetailView(LoginRequiredMixin, DetailView):
    """Document detail view"""
    model = Document
    template_name = 'rag_app/document_detail.html'
    context_object_name = 'document'
    
    def get_queryset(self):
        return Document.objects.filter(uploaded_by=self.request.user)
    
    def get_context_data(self, **kwargs):
        context = super().get_context_data(**kwargs)
        document = self.get_object()
        context['chunks'] = document.chunks.all()[:10]  # Show first 10 chunks
        context['total_chunks'] = document.chunks.count()
        return context


class DocumentUploadView(LoginRequiredMixin, CreateView):
    """Upload new document"""
    model = Document
    form_class = DocumentUploadForm
    template_name = 'rag_app/document_upload.html'
    success_url = reverse_lazy('rag_app:document_list')
    
    def get_form_kwargs(self):
        kwargs = super().get_form_kwargs()
        kwargs['user'] = self.request.user
        return kwargs
    
    def get_context_data(self, **kwargs):
        context = super().get_context_data(**kwargs)
        context['user_subjects'] = Subject.objects.filter(created_by=self.request.user)
        return context
    
    def form_valid(self, form):
        form.instance.uploaded_by = self.request.user
        response = super().form_valid(form)
        
        # Process document in background (in production, use Celery)
        try:
            processor = DocumentProcessor()
            processor.process_document(self.object)
            messages.success(self.request, 'Document uploaded and processed successfully!')
        except Exception as e:
            logger.error(f"Error processing document {self.object.id}: {str(e)}")
            messages.warning(self.request, 'Document uploaded but processing failed. Please try again.')
        
        return response


class DocumentDeleteView(LoginRequiredMixin, DeleteView):
    """Delete document"""
    model = Document
    template_name = 'rag_app/document_confirm_delete.html'
    success_url = reverse_lazy('rag_app:document_list')
    
    def get_queryset(self):
        return Document.objects.filter(uploaded_by=self.request.user)
    
    def delete(self, request, *args, **kwargs):
        messages.success(request, 'Document deleted successfully!')
        return super().delete(request, *args, **kwargs)


@login_required
def process_document(request, pk):
    """Process document for RAG"""
    document = get_object_or_404(Document, pk=pk, uploaded_by=request.user)
    
    try:
        processor = DocumentProcessor()
        processor.process_document(document)
        messages.success(request, 'Document processed successfully!')
    except Exception as e:
        logger.error(f"Error processing document {pk}: {str(e)}")
        messages.error(request, 'Error processing document. Please try again.')
    
    return redirect('rag_app:document_detail', pk=pk)


# Chat Views
class ChatView(LoginRequiredMixin, View):
    """Chat interface view"""
    template_name = 'rag_app/chat.html'
    
    def get(self, request, session_id=None):
        """Handle GET requests - show chat interface"""
        user = request.user
        
        # Check if a subject parameter is provided for starting a new subject chat
        subject_id = request.GET.get('subject')
        
        # Get or create chat session
        if session_id:
            session = get_object_or_404(ChatSession, id=session_id, user=user)
        elif subject_id:
            # Create a new chat session for the specified subject
            try:
                subject = get_object_or_404(Subject, id=subject_id, created_by=user)
                
                # Check if the subject has any processed documents
                has_documents = Document.objects.filter(
                    subject=subject, 
                    processed=True
                ).exists()
                
                if not has_documents:
                    messages.warning(
                        request, 
                        f'No processed documents found in "{subject.name}". Please upload and process some documents first.'
                    )
                    return redirect('rag_app:subject_detail', pk=subject.id)
                
                session = ChatSession.objects.create(
                    user=user,
                    subject=subject,
                    title=f"Chat with {subject.name}",
                    chat_type='subject'
                )
                # Redirect to the new session to avoid confusion with URL parameters
                return redirect('rag_app:chat_session', session_id=session.id)
            except (Subject.DoesNotExist, ValueError):
                # If subject doesn't exist or invalid ID, fall back to general chat
                session = ChatSession.objects.filter(user=user).order_by('-last_activity').first()
                if not session:
                    session = ChatSession.objects.create(user=user, title="New Chat")
        else:
            session = ChatSession.objects.filter(user=user).order_by('-last_activity').first()
            if not session:
                session = ChatSession.objects.create(user=user, title="New Chat")
        
        context = {
            'session': session,
            'current_session': session,  # For template compatibility
            'messages': session.messages.all() if session else [],
            'chat_sessions': ChatSession.objects.filter(user=user).order_by('-last_activity')[:10],
            'recent_sessions': ChatSession.objects.filter(user=user).order_by('-last_activity')[:10],
            'subjects': Subject.objects.filter(created_by=user),
            'form': ChatMessageForm(),
        }
        return render(request, self.template_name, context)
    
    def post(self, request, session_id=None):
        """Handle POST requests - send message"""
        try:
            user = request.user
            message_text = request.POST.get('message', '').strip()
            
            if not message_text:
                return JsonResponse({'error': 'Message cannot be empty'}, status=400)
            
            # Get or create session
            if session_id:
                session = get_object_or_404(ChatSession, id=session_id, user=user)
            else:
                # Get the most recent session for this user or create a new one
                session = ChatSession.objects.filter(user=user).order_by('-last_activity').first()
                if not session:
                    session = ChatSession.objects.create(
                        user=user,
                        title=message_text[:50] + "..." if len(message_text) > 50 else message_text
                    )
                else:
                    # Update the session title if it's still "New Chat" and this is the first user message
                    if session.title == "New Chat" and not session.messages.filter(is_user=True).exists():
                        session.title = message_text[:50] + "..." if len(message_text) > 50 else message_text
                        session.save()
            
            # Save user message
            user_message = ChatMessage.objects.create(
                session=session,
                message=message_text,
                is_user=True
            )
            
            # Generate AI response using RAG pipeline
            start_time = timezone.now()
            
            try:
                # Import RAG model
                from .pipeline.model import get_rag_model
                
                # Get RAG model instance
                rag_model = get_rag_model()
                
                # Check if user has any documents before allowing chat
                user_has_documents = Document.objects.filter(uploaded_by=user).exists()
                user_has_subjects_with_docs = Subject.objects.filter(
                    created_by=user, documents__isnull=False
                ).exists()
                
                # Process query based on session type
                if session.chat_type == 'anonymous' and session.temp_document:
                    # Anonymous document chat
                    rag_result = rag_model.query_temp_document(
                        question=message_text,
                        temp_document=session.temp_document,
                        chat_session=session
                    )
                elif session.subject:
                    # Subject-based chat with all documents from the subject
                    subject_has_docs = Document.objects.filter(subject=session.subject).exists()
                    if not subject_has_docs:
                        ai_response = f"No documents have been uploaded to the '{session.subject.name}' subject yet. Please upload some documents to this subject before starting a chat."
                    else:
                        rag_result = rag_model.query(
                            question=message_text,
                            subject_id=session.subject.id,
                            chat_session=session,
                            retrieval_strategy='hybrid',
                            max_chunks=5
                        )
                elif user_has_documents or user_has_subjects_with_docs:
                    # General chat with user's documents
                    rag_result = rag_model.query(
                        question=message_text,
                        subject_id=None,
                        chat_session=session,
                        retrieval_strategy='hybrid',
                        max_chunks=5
                    )
                else:
                    # No documents available - provide helpful guidance
                    ai_response = """Hello! I'm your AI study assistant. To get started, you'll need to upload some documents first. Here's how:

1. **Create a Subject**: Go to the Subjects section and create a new subject for your study material
2. **Upload Documents**: Add PDF, Word, PowerPoint, or text files to your subject
3. **Start Chatting**: Once documents are uploaded and processed, you can ask me questions about them

Alternatively, you can use the "Chat with Document" feature to quickly upload a single document and start chatting about it immediately.

What would you like to do first?"""
                
                # Only process RAG result if we didn't set a custom response
                if 'ai_response' not in locals() and 'rag_result' in locals():
                    if rag_result['success']:
                        ai_response = rag_result['answer']
                    else:
                        ai_response = rag_result.get('answer', 'I apologize, but I encountered an error while processing your question.')
                        logger.warning(f"RAG query failed: {rag_result.get('error')}")
                    
            except Exception as e:
                logger.error(f"Error using RAG model: {e}")
                # Fallback to simple response
                ai_response = "I'm having trouble accessing the document knowledge base right now. Please make sure documents are uploaded and try again."
            response_time = (timezone.now() - start_time).total_seconds()
            
            # Save AI message
            ai_message = ChatMessage.objects.create(
                session=session,
                message=ai_response,
                is_user=False,
                response_time=response_time
            )
            
            # Update session activity
            session.last_activity = timezone.now()
            session.save()
            
            return JsonResponse({
                'success': True,
                'response': ai_response,
                'session_id': str(session.id),
                'ai_message': {
                    'id': str(ai_message.id),
                    'message': ai_message.message,
                    'timestamp': ai_message.timestamp.isoformat(),
                }
            })
            
        except Exception as e:
            logger.error(f"Error in chat POST: {str(e)}")
            return JsonResponse({'error': 'An error occurred while processing your message'}, status=500)


@login_required
@csrf_exempt
def send_message(request):
    """Send chat message via AJAX"""
    if request.method == 'POST':
        try:
            data = json.loads(request.body)
            message_text = data.get('message', '').strip()
            session_id = data.get('session_id')
            subject_id = data.get('subject_id')
            
            if not message_text:
                return JsonResponse({'error': 'Message cannot be empty'}, status=400)
            
            # Get or create session
            if session_id:
                session = get_object_or_404(ChatSession, id=session_id, user=request.user)
            else:
                # Get the most recent session for this user or create a new one
                session = ChatSession.objects.filter(user=request.user).order_by('-last_activity').first()
                if not session:
                    session = ChatSession.objects.create(
                        user=request.user,
                        title=message_text[:50] + "..." if len(message_text) > 50 else message_text,
                        subject_id=subject_id if subject_id else None
                    )
                else:
                    # Update the session title if it's still "New Chat" and this is the first user message
                    if session.title == "New Chat" and not session.messages.filter(is_user=True).exists():
                        session.title = message_text[:50] + "..." if len(message_text) > 50 else message_text
                        session.save()
            
            # Save user message
            user_message = ChatMessage.objects.create(
                session=session,
                message=message_text,
                is_user=True
            )
            
            # Generate AI response using RAG pipeline
            start_time = timezone.now()
            
            try:
                # Import RAG model
                from .pipeline.model import get_rag_model
                
                # Get RAG model instance
                rag_model = get_rag_model()
                
                # Check if user has any documents before allowing chat
                user_has_documents = Document.objects.filter(uploaded_by=request.user).exists()
                user_has_subjects_with_docs = Subject.objects.filter(
                    created_by=request.user, documents__isnull=False
                ).exists()
                
                # Process query based on session type
                if session.chat_type == 'anonymous' and session.temp_document:
                    # Anonymous document chat
                    rag_result = rag_model.query_temp_document(
                        question=message_text,
                        temp_document=session.temp_document,
                        chat_session=session
                    )
                elif session.subject:
                    # Subject-based chat with all documents from the subject
                    subject_has_docs = Document.objects.filter(subject=session.subject).exists()
                    if not subject_has_docs:
                        ai_response = f"No documents have been uploaded to the '{session.subject.name}' subject yet. Please upload some documents to this subject before starting a chat."
                    else:
                        rag_result = rag_model.query(
                            question=message_text,
                            subject_id=session.subject.id,
                            chat_session=session,
                            retrieval_strategy='hybrid',
                            max_chunks=5
                        )
                elif user_has_documents or user_has_subjects_with_docs:
                    # General chat with user's documents
                    rag_result = rag_model.query(
                        question=message_text,
                        subject_id=None,
                        chat_session=session,
                        retrieval_strategy='hybrid',
                        max_chunks=5
                    )
                else:
                    # No documents available - provide helpful guidance
                    ai_response = """Hello! I'm your AI study assistant. To get started, you'll need to upload some documents first. Here's how:

1. **Create a Subject**: Go to the Subjects section and create a new subject for your study material
2. **Upload Documents**: Add PDF, Word, PowerPoint, or text files to your subject
3. **Start Chatting**: Once documents are uploaded and processed, you can ask me questions about them

Alternatively, you can use the "Chat with Document" feature to quickly upload a single document and start chatting about it immediately.

What would you like to do first?"""
                
                # Only process RAG result if we didn't set a custom response
                if 'ai_response' not in locals() and 'rag_result' in locals():
                    if rag_result['success']:
                        ai_response = rag_result['answer']
                        
                        # Store relevant chunks for this message
                        if rag_result.get('sources'):
                            chunk_ids = [chunk['chunk_id'] for chunk in rag_result['sources']]
                            # Note: Will link chunks after saving the message
                    else:
                        ai_response = rag_result.get('answer', 'I apologize, but I encountered an error while processing your question.')
                        logger.warning(f"RAG query failed: {rag_result.get('error')}")
                    
            except Exception as e:
                logger.error(f"Error using RAG model: {e}")
                # Fallback to simple response
                ai_response = "I'm having trouble accessing the document knowledge base right now. Please make sure documents are uploaded and try again."
            
            response_time = (timezone.now() - start_time).total_seconds()
            
            # Save AI message
            ai_message = ChatMessage.objects.create(
                session=session,
                message=ai_response,
                is_user=False,
                response_time=response_time
            )
            
            # Link relevant chunks if available
            try:
                if 'rag_result' in locals() and rag_result.get('success') and rag_result.get('sources'):
                    chunk_ids = [chunk['chunk_id'] for chunk in rag_result['sources']]
                    chunks = DocumentChunk.objects.filter(id__in=chunk_ids)
                    ai_message.relevant_chunks.set(chunks)
            except Exception as e:
                logger.warning(f"Error linking chunks to message: {e}")
            
            # Update session activity
            session.last_activity = timezone.now()
            session.save()
            
            # Prepare response data
            response_data = {
                'success': True,
                'session_id': str(session.id),
                'user_message': {
                    'id': str(user_message.id),
                    'message': user_message.message,
                    'timestamp': user_message.timestamp.isoformat()
                },
                'ai_message': {
                    'id': str(ai_message.id),
                    'message': ai_message.message,
                    'timestamp': ai_message.timestamp.isoformat(),
                    'response_time': response_time
                }
            }
            
            # Add source information if available
            if 'rag_result' in locals() and rag_result.get('success') and rag_result.get('sources'):
                response_data['sources'] = [
                    {
                        'document_title': chunk['document_title'],
                        'document_type': chunk['document_type'],
                        'page_number': chunk['page_number'],
                        'relevance_score': round(chunk['score'], 3)
                    }
                    for chunk in rag_result['sources'][:3]  # Limit to top 3 sources
                ]
            
            return JsonResponse(response_data)
            
        except Exception as e:
            logger.error(f"Error in chat: {str(e)}")
            return JsonResponse({'error': 'An error occurred while processing your message'}, status=500)
    
    return JsonResponse({'error': 'Invalid request method'}, status=405)


class AnonymousDocumentChatView(LoginRequiredMixin, View):
    """Anonymous document chat - upload a document and chat about it"""
    template_name = 'rag_app/anonymous_chat.html'
    
    def get(self, request):
        """Show anonymous chat upload form"""
        return render(request, self.template_name)
    
    def post(self, request):
        """Handle document upload and start chat session"""
        try:
            file = request.FILES.get('file')
            initial_question = request.POST.get('initial_question', '').strip()
            
            if not file:
                messages.error(request, 'Please select a file to upload.')
                return render(request, self.template_name)
            
            # Validate file
            if file.size > 50 * 1024 * 1024:  # 50MB limit
                messages.error(request, 'File size must be less than 50MB.')
                return render(request, self.template_name)
            
            # Create temporary document
            from django.utils import timezone
            temp_doc = TempDocument.objects.create(
                title=file.name.rsplit('.', 1)[0],  # Remove extension
                file=file,
                uploaded_by=request.user,
                expires_at=timezone.now() + timezone.timedelta(hours=24)
            )
            
            # Process document immediately for chat
            try:
                from .pipeline.data_processor import DocumentProcessor
                processor = DocumentProcessor()
                
                # Process the temporary document (adapt processor for temp docs)
                processor.process_temp_document(temp_doc)
                temp_doc.processed = True
                temp_doc.save()
                
            except Exception as e:
                logger.error(f"Error processing temp document {temp_doc.id}: {str(e)}")
                messages.error(request, 'Error processing document. Please try again.')
                return render(request, self.template_name)
            
            # Create chat session
            session_title = f"Chat about {temp_doc.title}"
            if initial_question:
                session_title = initial_question[:50] + "..." if len(initial_question) > 50 else initial_question
            
            chat_session = ChatSession.objects.create(
                user=request.user,
                temp_document=temp_doc,
                title=session_title,
                chat_type='anonymous'
            )
            
            # If there's an initial question, process it
            if initial_question:
                try:
                    # Save user message
                    user_message = ChatMessage.objects.create(
                        session=chat_session,
                        message=initial_question,
                        is_user=True
                    )
                    
                    # Generate AI response
                    from .pipeline.model import get_rag_model
                    rag_model = get_rag_model()
                    
                    # Process query with temp document
                    rag_result = rag_model.query_temp_document(
                        question=initial_question,
                        temp_document=temp_doc,
                        chat_session=chat_session
                    )
                    
                    if rag_result['success']:
                        ai_response = rag_result['answer']
                    else:
                        ai_response = "I've processed your document. What would you like to know about it?"
                    
                    # Save AI message
                    ChatMessage.objects.create(
                        session=chat_session,
                        message=ai_response,
                        is_user=False
                    )
                    
                except Exception as e:
                    logger.error(f"Error processing initial question: {str(e)}")
                    # Continue without the initial response
            
            # Redirect to chat session
            messages.success(request, 'Document uploaded successfully! You can now chat about it.')
            return redirect('rag_app:chat_session', session_id=chat_session.id)
            
        except Exception as e:
            logger.error(f"Error in anonymous chat upload: {str(e)}")
            messages.error(request, 'An error occurred while processing your document.')
            return render(request, self.template_name)


@login_required
def new_chat_session(request):
    """Create new chat session"""
    if request.method == 'POST':
        subject_id = request.POST.get('subject_id')
        session = ChatSession.objects.create(
            user=request.user,
            title="New Chat",
            subject_id=subject_id if subject_id else None
        )
        return redirect('rag_app:chat_session', session_id=session.id)
    
    return redirect('rag_app:chat')


@login_required
@csrf_exempt
def chat_with_subject(request):
    """Chat with documents from a specific subject"""
    if request.method == 'POST':
        try:
            data = json.loads(request.body)
            message_text = data.get('message', '').strip()
            subject_id = data.get('subject_id')
            session_id = data.get('session_id')
            
            if not message_text:
                return JsonResponse({'error': 'Message cannot be empty'}, status=400)
            
            if not subject_id:
                return JsonResponse({'error': 'Subject ID is required'}, status=400)
            
            # Verify user has access to the subject
            try:
                subject = Subject.objects.get(id=subject_id, created_by=request.user)
            except Subject.DoesNotExist:
                return JsonResponse({'error': 'Subject not found or access denied'}, status=403)
            
            # Get or create session
            if session_id:
                try:
                    session = ChatSession.objects.get(id=session_id, user=request.user, subject_id=subject_id)
                except ChatSession.DoesNotExist:
                    return JsonResponse({'error': 'Chat session not found'}, status=404)
            else:
                session = ChatSession.objects.create(
                    user=request.user,
                    subject=subject,
                    title=f"{subject.name}: {message_text[:30]}..." if len(message_text) > 30 else f"{subject.name}: {message_text}"
                )
            
            # Save user message
            user_message = ChatMessage.objects.create(
                session=session,
                message=message_text,
                is_user=True
            )
            
            # Generate AI response using RAG pipeline with subject filtering
            start_time = timezone.now()
            
            try:
                from .pipeline.model import get_rag_model
                
                rag_model = get_rag_model()
                
                # Use chat_with_subject method for better subject integration
                rag_result = rag_model.chat_with_subject(
                    question=message_text,
                    subject_id=subject_id,
                    chat_session=session
                )
                
                if rag_result['success']:
                    ai_response = rag_result['answer']
                    sources = rag_result.get('sources', [])
                else:
                    ai_response = rag_result.get('answer', f'I apologize, but I couldn\'t find relevant information in the {subject.name} documents to answer your question.')
                    sources = []
                    logger.warning(f"RAG query failed for subject {subject_id}: {rag_result.get('error')}")
                    
            except Exception as e:
                logger.error(f"Error using RAG model for subject {subject_id}: {e}")
                ai_response = f"I'm having trouble accessing the {subject.name} documents right now. Please make sure documents are uploaded for this subject and try again."
                sources = []
            
            response_time = (timezone.now() - start_time).total_seconds()
            
            # Save AI message
            ai_message = ChatMessage.objects.create(
                session=session,
                message=ai_response,
                is_user=False,
                response_time=response_time
            )
            
            # Link relevant chunks if available
            try:
                if sources:
                    chunk_ids = [chunk['chunk_id'] for chunk in sources]
                    chunks = DocumentChunk.objects.filter(id__in=chunk_ids)
                    ai_message.relevant_chunks.set(chunks)
            except Exception as e:
                logger.warning(f"Error linking chunks to message: {e}")
            
            # Update session activity
            session.last_activity = timezone.now()
            session.save()
            
            # Prepare response
            response_data = {
                'success': True,
                'session_id': str(session.id),
                'subject': {
                    'id': subject.id,
                    'name': subject.name,
                    'code': subject.code
                },
                'user_message': {
                    'id': str(user_message.id),
                    'message': user_message.message,
                    'timestamp': user_message.timestamp.isoformat()
                },
                'ai_message': {
                    'id': str(ai_message.id),
                    'message': ai_message.message,
                    'timestamp': ai_message.timestamp.isoformat(),
                    'response_time': response_time
                }
            }
            
            # Add source information
            if sources:
                response_data['sources'] = [
                    {
                        'document_title': chunk['document_title'],
                        'document_type': chunk['document_type'],
                        'page_number': chunk['page_number'],
                        'relevance_score': round(chunk['score'], 3)
                    }
                    for chunk in sources[:5]  # Limit to top 5 sources
                ]
                response_data['documents_used'] = len(set(chunk['document_id'] for chunk in sources))
            
            return JsonResponse(response_data)
            
        except Exception as e:
            logger.error(f"Error in subject chat: {str(e)}")
            return JsonResponse({'error': 'An error occurred while processing your message'}, status=500)
    
    return JsonResponse({'error': 'Invalid request method'}, status=405)


@login_required
def get_subject_documents(request, subject_id):
    """Get documents available for a subject"""
    try:
        subject = Subject.objects.get(id=subject_id, created_by=request.user)
        documents = Document.objects.filter(
            subject=subject,
            processed=True
        ).values('id', 'title', 'document_type', 'page_count', 'uploaded_at')
        
        return JsonResponse({
            'success': True,
            'subject': {
                'id': subject.id,
                'name': subject.name,
                'code': subject.code
            },
            'documents': list(documents),
            'total_documents': len(documents)
        })
        
    except Subject.DoesNotExist:
        return JsonResponse({'error': 'Subject not found'}, status=404)
    except Exception as e:
        logger.error(f"Error getting subject documents: {e}")
        return JsonResponse({'error': 'An error occurred'}, status=500)


# Quiz Views
class QuizListView(LoginRequiredMixin, ListView):
    """List all quizzes for the user"""
    model = Quiz
    template_name = 'rag_app/quiz_list.html'
    context_object_name = 'quizzes'
    paginate_by = 12
    
    def get_queryset(self):
        return Quiz.objects.filter(created_by=self.request.user).order_by('-created_at')


class QuizDetailView(LoginRequiredMixin, DetailView):
    """Quiz detail view"""
    model = Quiz
    template_name = 'rag_app/quiz_detail.html'
    context_object_name = 'quiz'
    
    def get_queryset(self):
        return Quiz.objects.filter(created_by=self.request.user)
    
    def get_context_data(self, **kwargs):
        context = super().get_context_data(**kwargs)
        quiz = self.get_object()
        context.update({
            'questions': quiz.questions.all(),
            'user_attempts': quiz.attempts.filter(user=self.request.user),
            'can_retake': True,  # Add logic for retake restrictions if needed
        })
        return context


class QuizCreateView(LoginRequiredMixin, CreateView):
    """Create new quiz"""
    model = Quiz
    form_class = QuizCreateForm
    template_name = 'rag_app/quiz_form.html'
    
    def get_form_kwargs(self):
        kwargs = super().get_form_kwargs()
        kwargs['user'] = self.request.user
        return kwargs
    
    def form_valid(self, form):
        form.instance.created_by = self.request.user
        response = super().form_valid(form)
        messages.success(self.request, 'Quiz created successfully!')
        return response
    
    def get_success_url(self):
        return reverse('rag_app:quiz_detail', kwargs={'pk': self.object.pk})


class QuizTakeView(LoginRequiredMixin, DetailView):
    """Take quiz view"""
    model = Quiz
    template_name = 'rag_app/quiz_take.html'
    context_object_name = 'quiz'
    
    def get_context_data(self, **kwargs):
        context = super().get_context_data(**kwargs)
        quiz = self.get_object()
        
        # Create or get existing attempt
        attempt, created = QuizAttempt.objects.get_or_create(
            quiz=quiz,
            user=self.request.user,
            is_completed=False,
            defaults={'total_points': sum(q.points for q in quiz.questions.all())}
        )
        
        context.update({
            'attempt': attempt,
            'questions': quiz.questions.all(),
            'existing_responses': {
                r.question_id: r for r in attempt.responses.all()
            } if not created else {}
        })
        return context


class QuizResultsView(LoginRequiredMixin, DetailView):
    """Quiz results view"""
    model = QuizAttempt
    template_name = 'rag_app/quiz_results.html'
    context_object_name = 'attempt'
    
    def get_queryset(self):
        return QuizAttempt.objects.filter(user=self.request.user, is_completed=True)


class QuizAttemptDetailView(LoginRequiredMixin, DetailView):
    """Quiz attempt detail view"""
    model = QuizAttempt
    template_name = 'rag_app/quiz_attempt_detail.html'
    context_object_name = 'attempt'
    
    def get_queryset(self):
        return QuizAttempt.objects.filter(user=self.request.user)


@login_required
def submit_quiz_attempt(request, pk):
    """Submit quiz attempt"""
    attempt = get_object_or_404(QuizAttempt, pk=pk, user=request.user, is_completed=False)
    
    if request.method == 'POST':
        try:
            # Process answers
            total_earned = 0
            for question in attempt.quiz.questions.all():
                answer_key = f'question_{question.id}'
                
                if question.question_type == 'mcq':
                    choice_id = request.POST.get(answer_key)
                    if choice_id:
                        choice = question.choices.get(id=choice_id)
                        is_correct = choice.is_correct
                        points = question.points if is_correct else 0
                        
                        QuizResponse.objects.update_or_create(
                            attempt=attempt,
                            question=question,
                            defaults={
                                'selected_choice': choice,
                                'is_correct': is_correct,
                                'points_earned': points
                            }
                        )
                        total_earned += points
                
                elif question.question_type in ['sa', 'fb']:
                    text_answer = request.POST.get(answer_key, '').strip()
                    # Simple text matching for now (can be enhanced with NLP)
                    is_correct = False  # Implement text comparison logic
                    points = question.points if is_correct else 0
                    
                    QuizResponse.objects.update_or_create(
                        attempt=attempt,
                        question=question,
                        defaults={
                            'text_answer': text_answer,
                            'is_correct': is_correct,
                            'points_earned': points
                        }
                    )
                    total_earned += points
            
            # Update attempt
            attempt.earned_points = total_earned
            attempt.completed_at = timezone.now()
            attempt.time_taken = attempt.completed_at - attempt.started_at
            attempt.is_completed = True
            attempt.calculate_score()
            
            messages.success(request, f'Quiz completed! Your score: {attempt.score:.1f}%')
            return redirect('rag_app:quiz_results', pk=attempt.pk)
            
        except Exception as e:
            logger.error(f"Error submitting quiz attempt {pk}: {str(e)}")
            messages.error(request, 'Error submitting quiz. Please try again.')
    
    return redirect('rag_app:quiz_take', pk=attempt.quiz.pk)


@login_required
def generate_quiz_questions(request, pk):
    """Generate quiz questions automatically"""
    quiz = get_object_or_404(Quiz, pk=pk, created_by=request.user)
    
    try:
        from .pipeline.model import QuizGenerator
        generator = QuizGenerator()
        
        if quiz.based_on_document:
            questions = generator.generate_questions_from_document(
                quiz.based_on_document, 
                num_questions=quiz.total_questions
            )
        else:
            # Generate from all documents in subject
            documents = quiz.subject.documents.filter(processed=True)
            questions = generator.generate_questions_from_documents(
                list(documents), 
                num_questions=quiz.total_questions
            )
        
        # Save generated questions
        for i, q_data in enumerate(questions, 1):
            question = Question.objects.create(
                quiz=quiz,
                question_text=q_data['question'],
                question_type=q_data['type'],
                explanation=q_data.get('explanation', ''),
                order=i
            )
            
            # Add choices for MCQ
            if q_data['type'] == 'mcq' and 'choices' in q_data:
                for j, choice_data in enumerate(q_data['choices'], 1):
                    question.choices.create(
                        choice_text=choice_data['text'],
                        is_correct=choice_data['is_correct'],
                        order=j
                    )
        
        messages.success(request, f'Generated {len(questions)} questions successfully!')
        
    except Exception as e:
        logger.error(f"Error generating quiz questions for {pk}: {str(e)}")
        messages.error(request, 'Error generating questions. Please try again.')
    
    return redirect('rag_app:quiz_detail', pk=pk)


# Profile Views
class ProfileView(LoginRequiredMixin, DetailView):
    """User profile view"""
    model = UserProfile
    template_name = 'rag_app/profile.html'
    context_object_name = 'profile'
    
    def get_object(self):
        profile, created = UserProfile.objects.get_or_create(user=self.request.user)
        return profile


class ProfileEditView(LoginRequiredMixin, UpdateView):
    """Edit user profile"""
    model = UserProfile
    form_class = UserProfileForm
    template_name = 'rag_app/profile_edit.html'
    success_url = reverse_lazy('rag_app:profile')
    
    def get_object(self):
        profile, created = UserProfile.objects.get_or_create(user=self.request.user)
        return profile
    
    def form_valid(self, form):
        messages.success(self.request, 'Profile updated successfully!')
        return super().form_valid(form)


# Study Session Views
class StudySessionListView(LoginRequiredMixin, ListView):
    """List study sessions"""
    model = StudySession
    template_name = 'rag_app/study_session_list.html'
    context_object_name = 'sessions'
    paginate_by = 20
    
    def get_queryset(self):
        return StudySession.objects.filter(user=self.request.user)


class StudySessionDetailView(LoginRequiredMixin, DetailView):
    """Study session detail"""
    model = StudySession
    template_name = 'rag_app/study_session_detail.html'
    context_object_name = 'session'
    
    def get_queryset(self):
        return StudySession.objects.filter(user=self.request.user)


# Slide Generation View
class SlideGenerationView(LoginRequiredMixin, TemplateView):
    """Generate slides from documents"""
    template_name = 'rag_app/slide_generate.html'
    
    def get_context_data(self, **kwargs):
        context = super().get_context_data(**kwargs)
        context['subjects'] = Subject.objects.filter(created_by=self.request.user)
        context['documents'] = Document.objects.filter(
            uploaded_by=self.request.user, processed=True
        )
        return context
    
    def post(self, request, *args, **kwargs):
        try:
            document_id = request.POST.get('document_id')
            topic = request.POST.get('topic', '')
            slide_count = int(request.POST.get('slide_count', 10))
            
            document = get_object_or_404(
                Document, id=document_id, uploaded_by=request.user
            )
            
            from .pipeline.model import SlideGenerator
            generator = SlideGenerator()
            slides = generator.generate_slides(document, topic, slide_count)
            
            return JsonResponse({
                'success': True,
                'slides': slides
            })
            
        except Exception as e:
            logger.error(f"Error generating slides: {str(e)}")
            return JsonResponse({
                'success': False,
                'error': 'Error generating slides. Please try again.'
            })


class ProfileView(LoginRequiredMixin, View):
    """User profile view for viewing and updating profile"""
    
    def get(self, request):
        """Display user profile"""
        try:
            profile = request.user.userprofile
        except UserProfile.DoesNotExist:
            # Create profile if it doesn't exist
            profile = UserProfile.objects.create(user=request.user)
        
        form = UserProfileForm(instance=profile, user=request.user)
        
        context = {
            'form': form,
            'profile': profile,
            'user': request.user
        }
        
        return render(request, 'rag_app/profile.html', context)
    
    def post(self, request):
        """Update user profile"""
        try:
            profile = request.user.userprofile
        except UserProfile.DoesNotExist:
            profile = UserProfile.objects.create(user=request.user)
        
        form = UserProfileForm(
            request.POST, 
            request.FILES, 
            instance=profile, 
            user=request.user
        )
        
        if form.is_valid():
            form.save()
            messages.success(request, 'Profile updated successfully!')
            return redirect('rag_app:profile')
        else:
            messages.error(request, 'Please correct the errors below.')
        
        context = {
            'form': form,
            'profile': profile,
            'user': request.user
        }
        
        return render(request, 'rag_app/profile.html', context)


# Slide Generator Views
class SlideDownloadView(LoginRequiredMixin, View):
    """Secure view for downloading generated PowerPoint files"""
    
    def get(self, request, filename):
        """Serve the PowerPoint file for download"""
        try:
            # Security check: only allow files the user owns or created
            # Extract user_id from filename (format: title_userid_timestamp.pptx)
            filename_parts = filename.replace('.pptx', '').split('_')
            if len(filename_parts) < 2:
                return HttpResponseForbidden("Invalid file access")
            
            try:
                file_user_id = int(filename_parts[-2])  # Second to last part should be user_id
            except (ValueError, IndexError):
                return HttpResponseForbidden("Invalid file access")
            
            # Check if user can access this file
            if request.user.id != file_user_id and not request.user.is_superuser:
                return HttpResponseForbidden("You don't have permission to access this file")
            
            # Construct file path
            file_path = os.path.join(settings.MEDIA_ROOT, 'generated_slides', filename)
            
            # Check if file exists
            if not os.path.exists(file_path):
                raise Http404("File not found")
            
            # Serve the file
            try:
                with open(file_path, 'rb') as fh:
                    response = HttpResponse(fh.read(), content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation')
                    response['Content-Disposition'] = f'attachment; filename="{filename}"'
                    response['Content-Length'] = os.path.getsize(file_path)
                    return response
            except IOError:
                raise Http404("File not found")
                
        except Exception as e:
            logger.error(f"Error serving slide file {filename}: {str(e)}")
            return HttpResponseServerError("Error accessing file")


class SlideGeneratorView(LoginRequiredMixin, View):
    """View for generating PowerPoint slides from uploaded documents"""
    template_name = 'rag_app/slide_generate.html'
    
    def get(self, request):
        """Render the slide generator form"""
        # Check if RAG model/LLM is available
        try:
            from .pipeline.model import RAGModel
            # Try to initialize to check if API key is configured
            rag_model = RAGModel()
            ai_available = True
        except Exception:
            ai_available = False
        
        context = {
            'page_title': 'Generate Slides',
            'user': request.user,
            'ai_available': ai_available
        }
        return render(request, self.template_name, context)
    
    def post(self, request):
        """Process slide generation request"""
        try:
            # Get uploaded files
            uploaded_files = request.FILES.getlist('documents')
            if not uploaded_files:
                return JsonResponse({
                    'success': False, 
                    'error': 'No files uploaded'
                }, status=400)
            
            # Get form data
            slide_count = request.POST.get('slide_count', 'auto')
            custom_slide_count = request.POST.get('custom_slide_count')
            template = request.POST.get('template', 'professional')
            title = request.POST.get('title', '')
            language = request.POST.get('language', 'en')
            instructions = request.POST.get('instructions', '')
            
            # Validate slide count
            if slide_count == 'custom':
                try:
                    slide_count = int(custom_slide_count)
                    if slide_count < 1 or slide_count > 50:
                        raise ValueError("Slide count must be between 1 and 50")
                except (ValueError, TypeError):
                    return JsonResponse({
                        'success': False,
                        'error': 'Invalid custom slide count'
                    }, status=400)
            elif slide_count != 'auto':
                try:
                    slide_count = int(slide_count)
                except ValueError:
                    slide_count = 'auto'
            
            # Process files and generate slides
            processor = SlideProcessor()
            result = processor.generate_slides(
                files=uploaded_files,
                slide_count=slide_count,
                template=template,
                title=title,
                language=language,
                instructions=instructions,
                user=request.user
            )
            
            if result['success']:
                return JsonResponse({
                    'success': True,
                    'message': 'Slides generated successfully!',
                    'download_url': result['download_url'],
                    'file_name': result['file_name']
                })
            else:
                return JsonResponse({
                    'success': False,
                    'error': result['error']
                }, status=500)
                
        except Exception as e:
            logger.error(f"Error in slide generation: {str(e)}")
            return JsonResponse({
                'success': False,
                'error': 'An error occurred while generating slides'
            }, status=500)


class SlideProcessor:
    """Advanced helper class for processing documents and generating PowerPoint slides with existing RAG LLM"""
    
    def __init__(self):
        self.supported_formats = ['.pdf', '.doc', '.docx', '.txt', '.ppt', '.pptx']
        # Initialize the existing RAG model
        try:
            from .pipeline.model import RAGModel
            self.rag_model = RAGModel()
            self.llm_available = True
        except Exception as e:
            logger.warning(f"Could not initialize RAG model: {str(e)}")
            self.rag_model = None
            self.llm_available = False
    
    def generate_slides(self, files, slide_count, template, title, language, instructions, user):
        """
        Main method to generate PowerPoint slides from uploaded documents using existing RAG LLM
        """
        try:
            # Step 1: Validate and process uploaded files
            processed_content = self._process_uploaded_files(files)
            if not processed_content:
                return {'success': False, 'error': 'No valid content found in uploaded files'}
            
            # Step 2: Extract and structure content
            structured_content = self._extract_content_structure(processed_content)
            
            # Step 3: Generate slide content using existing RAG LLM
            if self.llm_available and self.rag_model:
                slide_content_text = self._generate_ai_slide_content_with_rag(
                    structured_content, slide_count, instructions, language, title
                )
            else:
                # Fallback to basic generation
                slide_content_text = self._generate_basic_slide_content(
                    structured_content, slide_count, instructions, language, title
                )
            
            # Step 4: Create PowerPoint presentation with advanced styling
            presentation_path = self._create_advanced_powerpoint(
                slide_content_text, template, title, user
            )
            
            # Step 5: Return success response with download URL
            from django.urls import reverse
            download_url = reverse('rag_app:slide_download', kwargs={'filename': presentation_path})
            
            return {
                'success': True,
                'download_url': download_url,
                'file_name': presentation_path  # Return the actual filename
            }
            
        except Exception as e:
            logger.error(f"Error in slide generation: {str(e)}")
            return {'success': False, 'error': str(e)}
    
    def _generate_ai_slide_content_with_rag(self, structured_content, slide_count, instructions, language, title):
        """Generate slide content using the existing RAG model LLM"""
        try:
            # Prepare the content for AI processing
            content_summary = structured_content['full_text'][:4000]  # Limit content length
            
            # Determine slide count
            if slide_count == 'auto':
                slide_count = min(max(3, len(structured_content['sections'])), 10)
            
            # Create the prompt for the LLM
            prompt = f"""
            Create exactly {slide_count} slides based on the following document content. 
            
            Document Content:
            {content_summary}
            
            Requirements:
            - Language: {language}
            - Presentation Title: {title or 'Document Analysis'}
            - Additional Instructions: {instructions}
            - Each slide should have a clear title and 3-5 bullet points
            - Make the content educational and well-structured
            - Focus on key concepts and important information
            
            Format each slide exactly like this:
            ### Slide Title Here
            • First bullet point
            • Second bullet point
            • Third bullet point
            
            Please create engaging, informative slides that capture the essence of the document.
            Start with a title slide, then create content slides, and end with a summary if appropriate.
            """
            
            # Use the existing RAG model's LLM method
            messages = [
                {
                    "role": "system", 
                    "content": "You are an expert educational content creator that creates well-structured, engaging presentation slides."
                },
                {
                    "role": "user", 
                    "content": prompt
                }
            ]
            
            # Call the existing LLM method
            response = self.rag_model._generate_llm_response(messages)
            
            if response['success']:
                return response['answer']
            else:
                logger.error(f"LLM generation failed: {response.get('error', 'Unknown error')}")
                # Fallback to basic generation
                return self._generate_basic_slide_content(structured_content, slide_count, instructions, language, title)
            
        except Exception as e:
            logger.error(f"Error in RAG slide generation: {str(e)}")
            # Fallback to basic generation
            return self._generate_basic_slide_content(structured_content, slide_count, instructions, language, title)
    
    def _generate_basic_slide_content(self, structured_content, slide_count, instructions, language, title):
        """Fallback method for generating slide content without AI"""
        content = structured_content['full_text']
        sections = structured_content['sections']
        key_topics = structured_content['key_topics']
        
        if slide_count == 'auto':
            slide_count = min(max(3, len(sections)), 10)
        
        slides_text = f"### {title or 'Document Analysis'}\n"
        slides_text += f"• Overview of key concepts\n"
        slides_text += f"• Based on {len(structured_content['sources'])} document(s)\n"
        slides_text += f"• Educational content analysis\n\n"
        
        # Generate content slides
        if sections:
            for i, section in enumerate(sections[:slide_count-1]):
                slides_text += f"### {section['title'][:50]}\n"
                bullet_points = self._generate_bullet_points(section['content'])
                for bullet in bullet_points[:4]:
                    slides_text += f"• {bullet}\n"
                slides_text += "\n"
        else:
            # Generate slides based on key topics
            topics_per_slide = max(1, len(key_topics) // (slide_count - 1))
            for i in range(0, min(len(key_topics), (slide_count - 1) * topics_per_slide), topics_per_slide):
                slide_topics = key_topics[i:i + topics_per_slide]
                slides_text += f"### Key Concepts: {', '.join(slide_topics[:2])}\n"
                for topic in slide_topics[:4]:
                    slides_text += f"• Understanding {topic.title()}\n"
                slides_text += "\n"
        
        return slides_text
    
    def _load_image_stream(self, path_or_url):
        """Load image from path or URL"""
        if not path_or_url:
            return None
        
        try:
            if str(path_or_url).startswith("http"):
                import requests
                resp = requests.get(path_or_url, timeout=15)
                resp.raise_for_status()
                return BytesIO(resp.content)
            elif os.path.exists(path_or_url):
                return open(path_or_url, "rb")
            else:
                return None
        except Exception as e:
            logger.warning(f"Could not load image {path_or_url}: {str(e)}")
            return None
    
    def _create_advanced_powerpoint(self, slide_content_text, template, title, user):
        """Create PowerPoint presentation with advanced styling, background, and logo"""
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.dml.color import RGBColor
            from io import BytesIO
            import os
            from django.conf import settings
            
            # Create presentation
            prs = Presentation()
            slide_width, slide_height = prs.slide_width, prs.slide_height
            slide_layout = prs.slide_layouts[6]  # blank layout for full control
            
            # Define template colors
            template_colors = self._get_template_colors(template)
            
            # Load background and logo from media/images
            # Try multiple possible filenames for flexibility
            bg_image_paths = [
                os.path.join(settings.MEDIA_ROOT, 'images', 'ppt_background.jpg'),
                os.path.join(settings.MEDIA_ROOT, 'images', 'ppt.jpg'),
                os.path.join(settings.MEDIA_ROOT, 'images', 'background.jpg')
            ]
            logo_paths = [
                os.path.join(settings.MEDIA_ROOT, 'images', 'logo.png'),
                os.path.join(settings.MEDIA_ROOT, 'images', 'logoejust.png'),
                os.path.join(settings.MEDIA_ROOT, 'images', 'logo.jpg')
            ]
            
            bg_stream = None
            for bg_path in bg_image_paths:
                bg_stream = self._load_image_stream(bg_path)
                if bg_stream:
                    break
                    
            logo_stream = None
            for logo_path in logo_paths:
                logo_stream = self._load_image_stream(logo_path)
                if logo_stream:
                    break
            
            # Split content into slides
            slides = slide_content_text.split("###")
            
            for i, slide in enumerate(slides):
                if not slide.strip():
                    continue
                    
                lines = slide.strip().split("\n")
                slide_title = lines[0].strip()
                body_lines = [l.strip() for l in lines[1:] if l.strip() and l.strip().startswith('•')]
                
                slide_obj = prs.slides.add_slide(slide_layout)
                
                # Add background image
                if bg_stream:
                    bg_stream.seek(0)
                    slide_obj.shapes.add_picture(bg_stream, 0, 0, width=slide_width, height=slide_height)
                else:
                    # Add colored background
                    background = slide_obj.background
                    fill = background.fill
                    fill.solid()
                    fill.fore_color.rgb = template_colors['background']
                
                # Add title
                title_box = slide_obj.shapes.add_textbox(
                    Inches(0.5), Inches(0.3), 
                    slide_width - Inches(1), Inches(1.2)
                )
                title_tf = title_box.text_frame
                title_tf.word_wrap = True
                p = title_tf.add_paragraph()
                p.text = slide_title
                run = p.runs[0]
                run.font.size = Pt(36)
                run.font.bold = True
                run.font.color.rgb = template_colors['title']
                
                # Add content
                if body_lines:
                    content_box = slide_obj.shapes.add_textbox(
                        Inches(0.8), Inches(1.8), 
                        slide_width - Inches(1.6), slide_height - Inches(2.5)
                    )
                    content_tf = content_box.text_frame
                    content_tf.word_wrap = True
                    
                    for line in body_lines:
                        p = content_tf.add_paragraph()
                        # Remove bullet symbol if present and add it back for consistency
                        clean_line = line.replace('•', '').strip()
                        p.text = clean_line
                        p.level = 0
                        run = p.runs[0]
                        run.font.size = Pt(20)
                        run.font.color.rgb = template_colors['content']
                
                # Add logo (bottom-right corner)
                if logo_stream:
                    logo_stream.seek(0)
                    logo_height = Inches(0.8)
                    slide_obj.shapes.add_picture(
                        logo_stream,
                        slide_width - Inches(1.5),  # x position
                        Inches(0.1),   # y position
                        height=logo_height
                    )
            
            # Save presentation
            filename = f"{title or 'Generated_Presentation'}_{user.id}_{timezone.now().strftime('%Y%m%d_%H%M%S')}.pptx"
            filepath = os.path.join(settings.MEDIA_ROOT, 'generated_slides', filename)
            
            # Ensure directory exists
            os.makedirs(os.path.dirname(filepath), exist_ok=True)
            
            prs.save(filepath)
            
            # Close streams
            if bg_stream:
                bg_stream.close()
            if logo_stream:
                logo_stream.close()
            
            return filename
            
        except Exception as e:
            logger.error(f"Error creating advanced PowerPoint: {str(e)}")
            raise
    
    def _get_template_colors(self, template):
        """Get color scheme based on template"""
        try:
            from pptx.dml.color import RGBColor
        except ImportError:
            # Fallback if pptx is not available
            logger.warning("python-pptx not installed, using fallback colors")
            return {
                'background': (240, 248, 255),  # Light blue as tuple
                'title': (0, 102, 204),         # Blue as tuple
                'content': (51, 51, 51)         # Dark gray as tuple
            }
        
        color_schemes = {
            'professional': {
                'background': RGBColor(240, 248, 255),  # Light blue
                'title': RGBColor(0, 102, 204),         # Blue
                'content': RGBColor(51, 51, 51)         # Dark gray
            },
            'academic': {
                'background': RGBColor(248, 248, 255),  # Very light purple
                'title': RGBColor(75, 0, 130),          # Indigo
                'content': RGBColor(25, 25, 112)        # Navy blue
            },
            'creative': {
                'background': RGBColor(255, 248, 220),  # Light yellow
                'title': RGBColor(255, 140, 0),         # Orange
                'content': RGBColor(139, 69, 19)        # Brown
            },
            'minimal': {
                'background': RGBColor(255, 255, 255),  # White
                'title': RGBColor(64, 64, 64),          # Dark gray
                'content': RGBColor(96, 96, 96)         # Medium gray
            },
            'corporate': {
                'background': RGBColor(245, 245, 245),  # Light gray
                'title': RGBColor(0, 51, 102),          # Corporate blue
                'content': RGBColor(51, 51, 51)         # Dark gray
            }
        }
        
        return color_schemes.get(template, color_schemes['professional'])
    
    def _process_uploaded_files(self, files):
        """Process and extract text from uploaded files"""
        all_content = []
        
        for file in files:
            try:
                # Check file extension
                file_extension = file.name.lower().split('.')[-1]
                if f'.{file_extension}' not in self.supported_formats:
                    continue
                
                # Extract text based on file type
                if file_extension == 'pdf':
                    content = self._extract_pdf_content(file)
                elif file_extension in ['doc', 'docx']:
                    content = self._extract_word_content(file)
                elif file_extension == 'txt':
                    content = self._extract_text_content(file)
                elif file_extension in ['ppt', 'pptx']:
                    content = self._extract_powerpoint_content(file)
                else:
                    continue
                
                if content:
                    all_content.append({
                        'filename': file.name,
                        'content': content,
                        'type': file_extension
                    })
                    
            except Exception as e:
                logger.warning(f"Error processing file {file.name}: {str(e)}")
                continue
        
        return all_content
    
    def _extract_pdf_content(self, file):
        """Extract text from PDF file"""
        try:
            import PyPDF2
            import io
            
            pdf_reader = PyPDF2.PdfReader(io.BytesIO(file.read()))
            text = ""
            
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
            
            return text.strip()
        except Exception as e:
            logger.error(f"Error extracting PDF content: {str(e)}")
            return ""
    
    def _extract_word_content(self, file):
        """Extract text from Word document"""
        try:
            import docx
            import io
            
            doc = docx.Document(io.BytesIO(file.read()))
            text = ""
            
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            
            return text.strip()
        except Exception as e:
            logger.error(f"Error extracting Word content: {str(e)}")
            return ""
    
    def _extract_text_content(self, file):
        """Extract text from plain text file"""
        try:
            return file.read().decode('utf-8')
        except Exception as e:
            logger.error(f"Error extracting text content: {str(e)}")
            return ""
    
    def _extract_powerpoint_content(self, file):
        """Extract text from PowerPoint file"""
        try:
            from pptx import Presentation
            import io
            
            prs = Presentation(io.BytesIO(file.read()))
            text = ""
            
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            
            return text.strip()
        except Exception as e:
            logger.error(f"Error extracting PowerPoint content: {str(e)}")
            return ""
    
    def _extract_content_structure(self, processed_content):
        """Analyze and structure the extracted content"""
        combined_text = ""
        sources = []
        
        for content_item in processed_content:
            combined_text += f"\n--- From {content_item['filename']} ---\n"
            combined_text += content_item['content']
            sources.append(content_item['filename'])
        
        # Basic content analysis
        sections = self._identify_sections(combined_text)
        key_topics = self._extract_key_topics(combined_text)
        
        return {
            'full_text': combined_text,
            'sources': sources,
            'sections': sections,
            'key_topics': key_topics,
            'word_count': len(combined_text.split())
        }
    
    def _identify_sections(self, text):
        """Identify main sections in the text"""
        # Simple section identification based on common patterns
        import re
        
        sections = []
        lines = text.split('\n')
        
        for i, line in enumerate(lines):
            line = line.strip()
            if not line:
                continue
                
            # Look for headings (all caps, numbered, or specific patterns)
            if (re.match(r'^\d+\.', line) or 
                re.match(r'^[A-Z\s]{10,}$', line) or
                re.match(r'^(Chapter|Section|Part)', line, re.IGNORECASE)):
                sections.append({
                    'title': line,
                    'line_number': i,
                    'content': self._get_section_content(lines, i)
                })
        
        return sections[:10]  # Limit to first 10 sections
    
    def _get_section_content(self, lines, start_line):
        """Get content for a specific section"""
        content = []
        for i in range(start_line + 1, min(start_line + 20, len(lines))):
            if lines[i].strip():
                content.append(lines[i].strip())
            if len(content) > 10:  # Limit content length
                break
        return ' '.join(content)
    
    def _extract_key_topics(self, text):
        """Extract key topics from the text"""
        # Simple keyword extraction (in a real implementation, you'd use NLP)
        import re
        from collections import Counter
        
        # Remove common words and extract meaningful terms
        words = re.findall(r'\b[A-Za-z]{4,}\b', text.lower())
        common_words = {'that', 'this', 'with', 'have', 'will', 'from', 'they', 'been', 'said', 'each', 'which', 'their', 'time', 'more', 'very', 'when', 'come', 'may', 'into', 'over', 'think', 'also', 'your', 'work', 'life', 'only', 'can', 'still', 'should', 'after', 'being', 'now', 'made', 'before', 'here', 'through', 'when', 'where', 'much', 'take', 'than', 'only', 'think', 'know', 'just', 'first', 'could', 'right', 'would', 'about', 'there', 'what', 'some'}
        
        filtered_words = [word for word in words if word not in common_words and len(word) > 4]
        
        return [word for word, count in Counter(filtered_words).most_common(15)]
    
    def _generate_slide_content(self, structured_content, slide_count, instructions, language):
        """Generate content for slides using AI (simplified version)"""
        # This is a simplified version. In production, you'd use OpenAI or other AI services
        
        slides = []
        content = structured_content['full_text']
        sections = structured_content['sections']
        key_topics = structured_content['key_topics']
        
        # Determine actual slide count
        if slide_count == 'auto':
            slide_count = min(max(3, len(sections)), 15)
        
        # Generate title slide
        slides.append({
            'type': 'title',
            'title': self._generate_presentation_title(content, key_topics),
            'subtitle': f"Based on {len(structured_content['sources'])} document(s)",
            'content': []
        })
        
        # Generate content slides
        if sections:
            # Use identified sections
            for i, section in enumerate(sections[:slide_count-2]):
                slides.append({
                    'type': 'content',
                    'title': section['title'][:60],  # Limit title length
                    'content': self._generate_bullet_points(section['content'])
                })
        else:
            # Generate slides based on key topics
            topics_per_slide = max(1, len(key_topics) // (slide_count - 2))
            for i in range(0, min(len(key_topics), (slide_count - 2) * topics_per_slide), topics_per_slide):
                slide_topics = key_topics[i:i + topics_per_slide]
                slides.append({
                    'type': 'content',
                    'title': f"Key Concepts: {', '.join(slide_topics[:2])}",
                    'content': self._generate_bullet_points_from_topics(content, slide_topics)
                })
        
        # Generate conclusion slide
        if len(slides) < slide_count:
            slides.append({
                'type': 'conclusion',
                'title': 'Summary',
                'content': self._generate_summary_points(key_topics[:5])
            })
        
        return slides[:slide_count]
    
    def _generate_presentation_title(self, content, key_topics):
        """Generate a title for the presentation"""
        if key_topics:
            return f"{key_topics[0].title()} and Related Concepts"
        else:
            return "Document Analysis Presentation"
    
    def _generate_bullet_points(self, content):
        """Generate bullet points from content"""
        sentences = content.split('.')
        bullets = []
        
        for sentence in sentences[:5]:  # Max 5 bullet points
            sentence = sentence.strip()
            if len(sentence) > 20 and len(sentence) < 150:
                bullets.append(sentence.capitalize())
        
        return bullets or ["Key information from the document"]
    
    def _generate_bullet_points_from_topics(self, content, topics):
        """Generate bullet points based on topics"""
        bullets = []
        for topic in topics[:4]:  # Max 4 bullets
            bullets.append(f"Understanding {topic.title()}")
            bullets.append(f"Applications of {topic.title()}")
        
        return bullets[:5]  # Limit to 5 bullets
    
    def _generate_summary_points(self, key_topics):
        """Generate summary points"""
        return [f"Covered {topic.title()}" for topic in key_topics[:4]]
    
    def _create_powerpoint_presentation(self, slide_content, template, title, user):
        """Create the actual PowerPoint presentation"""
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.dml.color import RGBColor
            import os
            from django.conf import settings
            
            # Create presentation
            prs = Presentation()
            
            # Apply template styling
            self._apply_template_styling(prs, template)
            
            # Create slides
            for slide_data in slide_content:
                if slide_data['type'] == 'title':
                    self._create_title_slide(prs, slide_data, template)
                elif slide_data['type'] in ['content', 'conclusion']:
                    self._create_content_slide(prs, slide_data, template)
            
            # Save presentation
            filename = f"{title or 'Generated_Presentation'}_{user.id}_{timezone.now().strftime('%Y%m%d_%H%M%S')}.pptx"
            filepath = os.path.join(settings.MEDIA_ROOT, 'generated_slides', filename)
            
            # Ensure directory exists
            os.makedirs(os.path.dirname(filepath), exist_ok=True)
            
            prs.save(filepath)
            
            return filename
            
        except Exception as e:
            logger.error(f"Error creating PowerPoint: {str(e)}")
            raise
    
    def _apply_template_styling(self, prs, template):
        """Apply template-specific styling to the presentation"""
        # Template styling would be implemented here
        # For now, we'll use default styling
        pass
    
    def _create_title_slide(self, prs, slide_data, template):
        """Create a title slide"""
        slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = slide_data['title']
        subtitle.text = slide_data['subtitle']
    
    def _create_content_slide(self, prs, slide_data, template):
        """Create a content slide with bullet points"""
        slide_layout = prs.slide_layouts[1]  # Content slide layout
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = slide_data['title']
        
        tf = content.text_frame
        tf.clear()
        
        for bullet_point in slide_data['content']:
            p = tf.add_paragraph()
            p.text = bullet_point
            p.level = 0
